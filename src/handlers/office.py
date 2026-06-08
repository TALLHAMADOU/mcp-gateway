"""Local office document generation + LibreOffice headless conversion.

Runs entirely on the Linux host, no account required. Produces MS-native
files (.docx/.xlsx/.pptx) via python-docx / openpyxl / python-pptx, and
converts/exports (e.g. to PDF) via `soffice --headless`.

Heavy libraries are imported lazily so a missing optional dependency only
fails the specific tool, never the whole gateway.
"""
from fastapi import APIRouter, HTTPException
from starlette.concurrency import run_in_threadpool
import os
import shutil
import subprocess

from .fs import BASE as FS_BASE, _resolve

router = APIRouter()

OUTPUT_DIR = os.path.realpath(os.environ.get('OFFICE_OUTPUT_DIR') or os.path.join(FS_BASE, 'output'))


def _out_path(filename: str, default_ext: str) -> str:
    name = os.path.basename(filename or f'document.{default_ext}')
    if not name.endswith(f'.{default_ext}'):
        name += f'.{default_ext}'
    os.makedirs(OUTPUT_DIR, exist_ok=True)
    return os.path.join(OUTPUT_DIR, name)


@router.get('/health')
async def health():
    return {
        'status': 'ok',
        'handler': 'office',
        'output_dir': OUTPUT_DIR,
        'libreoffice': bool(shutil.which('soffice') or shutil.which('libreoffice')),
    }


# --- DOCX -------------------------------------------------------------------
def make_docx(spec: dict) -> str:
    try:
        from docx import Document
    except ImportError:
        raise RuntimeError('python-docx not installed')
    doc = Document()
    for el in spec.get('elements', []):
        t = el.get('type', 'paragraph')
        text = el.get('text', '')
        if t == 'heading':
            doc.add_heading(text, level=int(el.get('level', 1)))
        elif t == 'bullet':
            doc.add_paragraph(text, style='List Bullet')
        else:
            doc.add_paragraph(text)
    path = _out_path(spec.get('filename', 'document.docx'), 'docx')
    doc.save(path)
    return path


# --- XLSX -------------------------------------------------------------------
def make_xlsx(spec: dict) -> str:
    try:
        from openpyxl import Workbook
    except ImportError:
        raise RuntimeError('openpyxl not installed')
    wb = Workbook()
    sheets = spec.get('sheets') or [{'name': 'Sheet1', 'rows': spec.get('rows', [])}]
    wb.remove(wb.active)
    for s in sheets:
        ws = wb.create_sheet(title=str(s.get('name', 'Sheet')))
        for row in s.get('rows', []):
            ws.append(list(row))
    path = _out_path(spec.get('filename', 'data.xlsx'), 'xlsx')
    wb.save(path)
    return path


# --- PPTX -------------------------------------------------------------------
def make_pptx(spec: dict) -> str:
    try:
        from pptx import Presentation
    except ImportError:
        raise RuntimeError('python-pptx not installed')
    prs = Presentation()
    for slide in spec.get('slides', []):
        layout = prs.slide_layouts[1]  # title + content
        s = prs.slides.add_slide(layout)
        s.shapes.title.text = slide.get('title', '')
        body = s.placeholders[1].text_frame
        bullets = slide.get('bullets')
        if bullets:
            body.text = str(bullets[0])
            for b in bullets[1:]:
                body.add_paragraph().text = str(b)
        else:
            body.text = slide.get('body', '')
    path = _out_path(spec.get('filename', 'deck.pptx'), 'pptx')
    prs.save(path)
    return path


# --- LibreOffice conversion -------------------------------------------------
def convert_file(source_path: str, to: str) -> str:
    src = _resolve(source_path)
    if not os.path.isfile(src):
        raise FileNotFoundError('source file not found')
    soffice = shutil.which('soffice') or shutil.which('libreoffice')
    if not soffice:
        raise RuntimeError('LibreOffice (soffice) not installed on host')
    os.makedirs(OUTPUT_DIR, exist_ok=True)
    subprocess.run(
        [soffice, '--headless', '--convert-to', to, '--outdir', OUTPUT_DIR, src],
        check=True, capture_output=True, timeout=180,
    )
    ext = to.split(':', 1)[0]
    return os.path.join(OUTPUT_DIR, os.path.splitext(os.path.basename(src))[0] + '.' + ext)


# --- routes -----------------------------------------------------------------
def _wrap(fn, payload):
    try:
        return {'ok': True, 'path': fn(payload)}
    except (RuntimeError, FileNotFoundError) as e:
        raise HTTPException(status_code=400, detail=str(e))
    except subprocess.CalledProcessError as e:
        raise HTTPException(status_code=500, detail=e.stderr.decode(errors='replace'))


@router.post('/docx')
async def create_docx(payload: dict):
    return await run_in_threadpool(_wrap, make_docx, payload)


@router.post('/xlsx')
async def create_xlsx(payload: dict):
    return await run_in_threadpool(_wrap, make_xlsx, payload)


@router.post('/pptx')
async def create_pptx(payload: dict):
    return await run_in_threadpool(_wrap, make_pptx, payload)


@router.post('/convert')
async def convert(payload: dict):
    src, to = payload.get('source_path'), payload.get('to')
    if not src or not to:
        raise HTTPException(status_code=400, detail='source_path and to are required')
    try:
        return {'ok': True, 'path': await run_in_threadpool(convert_file, src, to)}
    except (RuntimeError, FileNotFoundError) as e:
        raise HTTPException(status_code=400, detail=str(e))
    except subprocess.CalledProcessError as e:
        raise HTTPException(status_code=500, detail=e.stderr.decode(errors='replace'))
